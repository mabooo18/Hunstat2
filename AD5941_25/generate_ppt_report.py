from pptx import Presentation
from pptx.util import Pt

output_path = r"c:\Users\ASUS\Downloads\potensio\HunStat2-main\HunStat2-main\Software\AD5941_25\Presentasi_Refaktorisasi_HunStat2.pptx"

prs = Presentation()


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)
        first = False


add_title_slide(
    "Refaktorisasi HunStat2 / AD5941_25",
    "Laporan implementasi, debugging, dan kesiapan presentasi"
)

add_bullet_slide("1. Latar Belakang", [
    "Kode awal cenderung monolitik dan sulit dipelihara",
    "Debugging data plot kosong memakan waktu",
    "Diperlukan struktur modular dan alat uji cepat"
])

add_bullet_slide("2. Tujuan Pekerjaan", [
    "Memodularisasi firmware agar mirip pola FreiStat",
    "Memisahkan setup, komunikasi, storage, dan metode",
    "Menyediakan UI test Python untuk board dan dummy",
    "Membuat jalur diagnosis register AD5941"
])

add_bullet_slide("3. Arsitektur Baru Firmware", [
    "Main sketch menjadi orkestrator alur",
    "Layer setup untuk init/reset/kalibrasi AD5941",
    "Layer communication untuk parse dan dispatch command",
    "Layer electrochemical_methods per metode (OCP/EIS/CA/SWV/DPV/CV)",
    "Layer interface/status untuk LED dan indikator"
])

add_bullet_slide("4. Manfaat Refraktorisasi", [
    "Kode lebih mudah dibaca dan diuji",
    "Risiko perubahan silang antar fitur berkurang",
    "Penambahan fitur baru lebih aman",
    "Troubleshooting lebih terarah"
])

add_bullet_slide("5. UI Python Test Console", [
    "Koneksi serial, kirim command, monitor log",
    "Atur parameter metode dari UI",
    "Plot data real-time per mode",
    "Export data ke CSV untuk analisis lanjutan"
])

add_bullet_slide("6. Mode Data: BOARD vs DUMMY", [
    "BOARD: command dikirim ke perangkat nyata",
    "DUMMY: data sintetis dibuat di UI untuk validasi alur",
    "Memudahkan demo tanpa ketergantungan hardware",
    "Mendukung uji parser dan plotting secara cepat"
])

add_bullet_slide("7. Investigasi Plot Kosong", [
    "UI hanya memplot data yang formatnya parseable",
    "Sebagian output firmware tidak selalu sesuai parser",
    "Solusi: perluasan parser + dummy mode untuk isolasi masalah"
])

add_bullet_slide("8. Investigasi Register AD5941", [
    "Dibuat test register terpisah (Python + Arduino sketch)",
    "Diverifikasi jalur command baca register tersedia",
    "CHIPID invalid (0x0000/0xFFFF) mengarah ke isu hardware/SPI",
    "Kesimpulan: bedakan masalah software vs wiring/power"
])

add_bullet_slide("9. Status Dummy Metode", [
    "EIS/CV/CA/SWV/DPV/OCP tersedia untuk pengujian UI",
    "Profil dummy bisa disetel sesuai kebutuhan demo",
    "SWV saat ini mengikuti preferensi presentasi pengguna"
])

add_bullet_slide("10. Hasil dan Dampak", [
    "Maintainability meningkat signifikan",
    "Debug loop lebih cepat (serial + plot dalam satu UI)",
    "Diagnosis AD5941 lebih sistematis",
    "Fondasi proyek lebih siap untuk pengembangan berikutnya"
])

add_bullet_slide("11. Rekomendasi Lanjutan", [
    "Tambah profil dummy realistis (linear/realistic/noisy)",
    "Standarisasi format output serial lintas metode",
    "Tambah metadata eksperimen dan timestamp di CSV",
    "Tambahkan test parser otomatis berbasis sample log"
])

add_bullet_slide("12. Penutup", [
    "Refaktorisasi berhasil memperkuat fondasi software",
    "UI test membantu validasi dan presentasi teknis",
    "Langkah selanjutnya: hardening build dan validasi board end-to-end"
])

prs.save(output_path)
print(output_path)
